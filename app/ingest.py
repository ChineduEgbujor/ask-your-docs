import os, glob, pickle
import numpy as np
from sentence_transformers import SentenceTransformer
import faiss
from dotenv import load_dotenv
import fitz  # PyMuPDF for PDF handling
import docx
from pptx import Presentation

DOCS_GLOB = "*.txt,*.pdf,*.docx,*.pptx".split(",")

load_dotenv()
model = SentenceTransformer("all-MiniLM-L6-v2")
INDEX_PATH = os.getenv("VECTOR_STORE_PATH")

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from a PDF file."""
    doc = fitz.open(pdf_path)
    return "\n".join(page.get_text() for page in doc)

def extract_text_from_docx(docx_path: str) -> str:
    """Extract text from a DOCX file."""
    doc = docx.Document(docx_path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from a PPTX file."""
    prs = Presentation(pptx_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text(path: str) -> str:
    """Extract text from various document types."""
    ext = path.rsplit(".", 1)[-1].lower()
    if ext == "txt":
        return open(path, encoding="utf-8").read()
    if ext == "pdf":
        return extract_text_from_pdf(path)
    if ext == "docx":
        return extract_text_from_docx(path)
    if ext == "pptx":
        return extract_text_from_pptx(path)
    # fallback: try to read as text
    return open(path, errors="ignore").read()

def ingest_docs(doc_dir: str):
    # 1) collect all matching files
    patterns = [os.path.join(doc_dir, p) for p in DOCS_GLOB]
    file_paths = []
    for pattern in patterns:
        file_paths.extend(glob.glob(pattern))
    texts, ids = [], []
    # 2) extract text
    for path in file_paths:
        text = extract_text(path)
        if not text.strip():
            continue
        texts.append(text)
        ids.append(os.path.basename(path))
    # 3) compute embeddings
    embs = model.encode(texts, show_progress_bar=True)
    embs = np.array(embs, dtype="float32")
    # 4) build FAISS index
    dim = embs.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(embs)
    os.makedirs(os.path.dirname(INDEX_PATH), exist_ok=True)
    faiss.write_index(index, INDEX_PATH)
    # 5) save metadata
    with open(INDEX_PATH + ".pkl", "wb") as f:
        pickle.dump({"ids": ids, "texts": texts}, f)
    print(f"Ingested {len(texts)} docs into {INDEX_PATH}")


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Ingest a directory of .txt files into a FAISS index"
    )
    parser.add_argument(
        "--doc-dir",
        required=True,
        help="Path to the folder containing your .txt documents"
    )
    args = parser.parse_args()

    # actually run the ingestion
    ingest_docs(args.doc_dir)